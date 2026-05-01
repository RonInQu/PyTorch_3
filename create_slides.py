"""Generate PowerPoint slides describing the Live Clot Detector system architecture."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MEDIUM_BLUE = RGBColor(0x2E, 0x5D, 0x8C)
LIGHT_BLUE = RGBColor(0x4A, 0x90, 0xD9)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Title background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    # Title text
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = LIGHT_BLUE
        p2.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(title, bullets, sub_bullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    # Body
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
        if sub_bullets and i in sub_bullets:
            for sb in sub_bullets[i]:
                p2 = tf2.add_paragraph()
                p2.text = "    " + sb
                p2.font.size = Pt(15)
                p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
                p2.space_after = Pt(4)
    return slide

def add_architecture_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture — Inference Pipeline"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Draw pipeline boxes
    box_y = Inches(2.0)
    box_h = Inches(1.2)
    box_w = Inches(2.0)
    arrow_w = Inches(0.4)
    
    stages = [
        ("R(t)\n150 Hz", LIGHT_GRAY, DARK_GRAY),
        ("5s Sliding\nWindow", LIGHT_BLUE, WHITE),
        ("Feature\nExtraction\n(21 feats)", MEDIUM_BLUE, WHITE),
        ("StandardScaler\n(z-score)", MEDIUM_BLUE, WHITE),
        ("GRU\nSequence\n(len=8)", ACCENT_RED, WHITE),
        ("Softmax/T\n(T=1.5)", MEDIUM_BLUE, WHITE),
        ("EMA\nSmoothing", MEDIUM_BLUE, WHITE),
        ("DA Override\n+ Safety", ACCENT_GREEN, WHITE),
        ("P(blood)\nP(clot)\nP(wall)", LIGHT_GRAY, DARK_GRAY),
    ]
    
    start_x = Inches(0.3)
    for i, (label, fill_color, text_color) in enumerate(stages):
        x = start_x + i * (box_w + arrow_w)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, box_y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
        box.line.color.rgb = DARK_BLUE
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        
        # Arrow between boxes
        if i < len(stages) - 1:
            arrow_x = x + box_w
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, 
                                           arrow_x, box_y + Inches(0.4), 
                                           arrow_w, Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = DARK_BLUE
            arrow.line.fill.background()

    # Annotations below
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12), Inches(3.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    annotations = [
        "• Input: Raw resistance at 150 Hz from catheter sensor",
        "• Window: 5 seconds (~750 samples), reported every 200 ms",
        "• Features: 21 selected from 64 candidates (clot_wall_focused set)",
        "• Model: GRU (hidden=32) → FC(24) → FC(3), sequence length = 8 windows",
        "• Post-processing: Temperature scaling → EMA posterior smoothing → DA override logic",
        "• Output: Real-time probability vector [P(blood), P(clot), P(wall)] every 200 ms",
    ]
    for i, ann in enumerate(annotations):
        p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
        p.text = ann
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    return slide

# ═══════════════════════════════════════════
# BUILD SLIDES
# ═══════════════════════════════════════════

# Slide 1: Title
add_title_slide("Live Clot Detector — System Architecture",
                "Real-Time GRU-Based Tissue Classification\nV6 — April 2026")

# Slide 2: Problem Statement
add_content_slide("Problem Statement", [
    "Goal: Real-time classification of catheter contact during procedure",
    "Three tissue states: Blood (baseline), Clot, Wall",
    "Current system: Detection Algorithm (DA) — rule-based, from LED state",
    "Challenge: DA misclassifies ~30% of clot/wall events",
    "ML Objective: Improve over DA by learning from resistance signal dynamics",
    "Constraint: Must run in real-time at 150 Hz input rate",
])

# Slide 3: Architecture
add_architecture_slide()

# Slide 4: Feature Extraction
add_content_slide("Feature Extraction (21 Active Features)", [
    "64 total features computed from 5-second sliding window",
    "21 selected via importance analysis (clot_wall_focused set):",
    "Basic Statistics: mean, std, min, max, range, windowed diff",
    "Derivative Stats: kurtosis, mean_abs, std, skew",
    "EMA Signals: slow EMA, abs_diff, ratio",
    "Detrended: std, std300, skew",
    "Percentiles: p95-p5, frac_above_p95",
    "Dynamics: Hjorth mobility/complexity, mean abs 2nd derivative",
], sub_bullets={
    1: ["Selected by Cohen's d (clot-vs-wall separation) + AUC ranking",
        "Feature indices: [39,21,4,19,41,9,5,23,0,34,28,29,3,38,17,32,42,27,1,20,40]"]
})

# Slide 5: GRU Model
add_content_slide("ClotGRU Model", [
    "Architecture: GRU(input=21, hidden=32) → FC(32→24, ReLU) → FC(24→3)",
    "Sequence length: 8 windows (8 × 200ms = 1.6 seconds of context)",
    "Total parameters: ~5,000 (lightweight for real-time)",
    "Training: CrossEntropyLoss, Adam (lr=1e-4, wd=1e-4), 100 epochs, patience=15",
    "Input normalization: StandardScaler fit on training data features",
    "Weight initialization: Orthogonal (GRU), Kaiming (FC1), Xavier (FC2)",
    "Output: 3-class logits → Temperature-scaled softmax (T=1.5)",
])

# Slide 6: Post-Processing
add_content_slide("Post-Processing Pipeline", [
    "1. Temperature Scaling (T=1.5): Softens overconfident predictions",
    "2. EMA Posterior Smoothing: Blends new probs with running posterior",
    "3. DA Override Logic: Device-assisted labels can override GRU",
    "4. Post-EMA Safety Net: Final check — force DA label if GRU not confident",
], sub_bullets={
    1: ["Blood→tissue: α_new=0.22 (moderate reactivity)",
        "Tissue→blood: α_new=0.65 (fast exit)",
        "Same class: α_new=0.03 (very stable)",
        "Cross class: α_new=0.01 (resist flicker)"],
    2: ["DA=blood → hard reset (clear GRU hidden state)",
        "DA=clot/wall → override unless GRU P(top) > threshold (0.80/0.92)"],
})

# Slide 7: Training Pipeline
add_content_slide("Training Pipeline", [
    "1. Labeling (Labeling_5Names_V6): Crop to first-last blood, blank artifacts",
    "2. Scaler Fitting (fit_scaler_V6): StandardScaler on windowed training features",
    "3. Model Training (train_gru_V6): GRU with cached feature extraction",
    "4. Inference (gru_torch_V6): Full pipeline on held-out test studies",
    "",
    "Data Preparation:",
    "   • All non-tissue events blanked to blood_median + noise(σ=5)",
    "   • Artifacts (contrast=event 8, saline=event 15) always blanked",
    "   • R > 5000 Ω outliers blanked in training, clipped in testing",
    "   • Window: 5s, stride: 30 samples (~200ms), sequence: 8 windows",
])

# Slide 8: Feature Scaling — How It Works
add_content_slide("Feature Scaling — StandardScaler (Z-Score Normalization)", [
    "Problem: Raw features have vastly different scales",
    "   • mean R ≈ 50–5000 Ω,  std R ≈ 1–500,  kurtosis ≈ -2 to +20",
    "   • GRU gradient updates dominated by large-magnitude features",
    "",
    "Solution: StandardScaler transforms each feature independently:",
    "        x_scaled = (x - μ) / σ",
    "   where μ = training mean, σ = training std for that feature",
    "",
    "After scaling: each feature has mean ≈ 0, std ≈ 1",
    "   • All features contribute equally to learning",
    "   • Scaler is fit ONLY on training data (no data leakage)",
    "   • Same μ, σ applied at inference time (saved as .pkl)",
])

# Slide 9: Feature Distributions Before & After Scaling
def add_distribution_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Feature Distributions — Before vs After Scaling"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # LEFT: Before scaling
    left_x = Inches(0.4)
    col_w = Inches(6.0)
    
    # Left header
    hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, Inches(1.5), col_w, Inches(0.6))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = ACCENT_RED
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "BEFORE Scaling (Raw Features)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Left table: example raw feature ranges
    raw_data = [
        ("Feature", "Range", "Mean", "Std"),
        ("mean (R)", "50 – 5000", "~800", "~600"),
        ("std (R)", "0.5 – 500", "~40", "~60"),
        ("kurtosis deriv", "-2 – 20", "~3", "~4"),
        ("frac_above_p95", "0 – 0.2", "~0.05", "~0.03"),
        ("Hjorth mobility", "0 – 50", "~5", "~8"),
        ("peak-to-peak", "1 – 4000", "~300", "~400"),
    ]
    
    row_h = Inches(0.38)
    table_top = Inches(2.3)
    for i, (feat, rng, mean, std) in enumerate(raw_data):
        y_pos = table_top + i * row_h
        # Background for header row
        if i == 0:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, y_pos, col_w, row_h)
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            bg.line.fill.background()
        
        col_widths = [Inches(2.0), Inches(1.5), Inches(1.2), Inches(1.2)]
        texts = [feat, rng, mean, std]
        cx = left_x
        for j, (txt, cw) in enumerate(zip(texts, col_widths)):
            tb = slide.shapes.add_textbox(cx, y_pos, cw, row_h)
            tf = tb.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = txt
            p.font.size = Pt(12)
            p.font.bold = (i == 0)
            p.font.color.rgb = DARK_GRAY
            cx += cw

    # Left annotation
    note_y = table_top + len(raw_data) * row_h + Inches(0.2)
    tb = slide.shapes.add_textbox(left_x, note_y, col_w, Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚠ Features span orders of magnitude"
    p.font.size = Pt(13)
    p.font.color.rgb = ACCENT_RED
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = "GRU weights skew toward high-magnitude features"
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GRAY

    # RIGHT: After scaling  
    right_x = Inches(6.9)
    
    # Right header
    hdr2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, Inches(1.5), col_w, Inches(0.6))
    hdr2.fill.solid()
    hdr2.fill.fore_color.rgb = ACCENT_GREEN
    hdr2.line.fill.background()
    tf = hdr2.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "AFTER Scaling (Z-Score)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Right table: after scaling
    scaled_data = [
        ("Feature", "Range", "Mean", "Std"),
        ("mean (R)", "-1.3 – 7.0", "0.0", "1.0"),
        ("std (R)", "-0.7 – 7.7", "0.0", "1.0"),
        ("kurtosis deriv", "-1.3 – 4.3", "0.0", "1.0"),
        ("frac_above_p95", "-1.7 – 5.0", "0.0", "1.0"),
        ("Hjorth mobility", "-0.6 – 5.6", "0.0", "1.0"),
        ("peak-to-peak", "-0.8 – 9.3", "0.0", "1.0"),
    ]
    
    for i, (feat, rng, mean, std) in enumerate(scaled_data):
        y_pos = table_top + i * row_h
        if i == 0:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, y_pos, col_w, row_h)
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            bg.line.fill.background()
        
        col_widths = [Inches(2.0), Inches(1.5), Inches(1.2), Inches(1.2)]
        texts = [feat, rng, mean, std]
        cx = right_x
        for j, (txt, cw) in enumerate(zip(texts, col_widths)):
            tb = slide.shapes.add_textbox(cx, y_pos, cw, row_h)
            tf = tb.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = txt
            p.font.size = Pt(12)
            p.font.bold = (i == 0)
            p.font.color.rgb = DARK_GRAY
            cx += cw

    # Right annotation
    tb2 = slide.shapes.add_textbox(right_x, note_y, col_w, Inches(1.2))
    tf = tb2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "✓ All features centered at 0, unit variance"
    p.font.size = Pt(13)
    p.font.color.rgb = ACCENT_GREEN
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = "Equal contribution to gradient updates"
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GRAY

    # Bottom: Arrow in the middle showing transformation
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, 
                                   Inches(6.1), Inches(3.5), Inches(0.7), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = DARK_BLUE
    arrow.line.fill.background()
    
    # Formula box at bottom center
    formula_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                          Inches(4.0), Inches(6.2), Inches(5.3), Inches(0.9))
    formula_box.fill.solid()
    formula_box.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    formula_box.line.color.rgb = MEDIUM_BLUE
    formula_box.line.width = Pt(2)
    tf = formula_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "x_scaled = (x − μ_train) / σ_train    →    Same transform at inference"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    return slide

add_distribution_slide()

# ═══════════════════════════════════════════
# SLIDES 10-12: GRU TEMPORAL ARCHITECTURE DETAIL
# ═══════════════════════════════════════════

# Slide 10: The 5-Second Window — What the GRU Sees
def add_window_detail_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "The 5-Second Window — Feature Extraction Unit"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Visual: timeline showing one 5s window
    # Draw a long horizontal bar representing the resistance signal
    signal_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                         Inches(1.0), Inches(2.0), Inches(11.0), Inches(0.6))
    signal_bar.fill.solid()
    signal_bar.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    signal_bar.line.color.rgb = DARK_GRAY
    signal_bar.line.width = Pt(1)
    # Label it
    tb = slide.shapes.add_textbox(Inches(4.5), Inches(1.5), Inches(4), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Continuous R(t) signal at 150 Hz"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Highlight one 5s window
    window_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                         Inches(5.5), Inches(1.9), Inches(3.5), Inches(0.8))
    window_box.fill.solid()
    window_box.fill.fore_color.rgb = RGBColor(0x4A, 0x90, 0xD9)  # LIGHT_BLUE
    window_box.fill.fore_color.brightness = 0.6
    window_box.line.color.rgb = LIGHT_BLUE
    window_box.line.width = Pt(2)
    # "5 sec" label
    tb2 = slide.shapes.add_textbox(Inches(6.2), Inches(2.0), Inches(2.0), Inches(0.5))
    tf = tb2.text_frame
    p = tf.paragraphs[0]
    p.text = "5.0 sec = 750 samples"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Arrow down to feature vector
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, 
                                    Inches(7.0), Inches(2.9), Inches(0.5), Inches(0.7))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = DARK_BLUE
    arrow.line.fill.background()

    # Feature vector box
    feat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                       Inches(4.5), Inches(3.8), Inches(5.5), Inches(0.7))
    feat_box.fill.solid()
    feat_box.fill.fore_color.rgb = MEDIUM_BLUE
    feat_box.line.color.rgb = DARK_BLUE
    feat_box.line.width = Pt(2)
    tf = feat_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Feature Vector: 21 values (scaled)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Explanation bullets below
    txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.5), Inches(2.5))
    tf = txBox3.text_frame
    tf.word_wrap = True
    explanations = [
        "• Each window captures 750 resistance samples (5s × 150 Hz)",
        "• From this window, 21 statistical features are computed:",
        "    mean, std, min, max, range, kurtosis, derivatives, EMA signals, detrended stats, etc.",
        "• The window slides by 30 samples (200 ms) → new feature vector every 200 ms",
        "• One window = one 'snapshot' of the signal's statistical properties",
        "• Why 5 seconds? Long enough to capture clot/wall dynamics, short enough for real-time",
    ]
    for i, txt in enumerate(explanations):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = txt
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)
    return slide

add_window_detail_slide()

# Slide 11: Sequence of 8 — Temporal Context
def add_sequence_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Sequence Length = 8 — How the GRU Sees Temporal Context"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Draw 8 boxes in a row representing the sequence
    box_w = Inches(1.3)
    box_h = Inches(1.5)
    gap = Inches(0.15)
    start_x = Inches(0.8)
    top_y = Inches(1.8)
    
    time_labels = ["t-7", "t-6", "t-5", "t-4", "t-3", "t-2", "t-1", "t (now)"]
    sec_labels = ["-1.4s", "-1.2s", "-1.0s", "-0.8s", "-0.6s", "-0.4s", "-0.2s", "0.0s"]
    
    for i in range(8):
        x = start_x + i * (box_w + gap)
        # Gradient effect: older = lighter, newer = darker
        alpha = 0.4 + 0.6 * (i / 7.0)
        r = int(0x2E * alpha + 0xE0 * (1 - alpha))
        g = int(0x5D * alpha + 0xE0 * (1 - alpha))
        b = int(0x8C * alpha + 0xE0 * (1 - alpha))
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(r, g, b)
        box.line.color.rgb = DARK_BLUE
        box.line.width = Pt(1.5)
        
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = time_labels[i]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE if i >= 4 else DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = "21 feats"
        p2.font.size = Pt(10)
        p2.font.color.rgb = WHITE if i >= 4 else DARK_GRAY
        p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph()
        p3.text = sec_labels[i]
        p3.font.size = Pt(10)
        p3.font.color.rgb = WHITE if i >= 4 else DARK_GRAY
        p3.alignment = PP_ALIGN.CENTER

    # Arrow from sequence to GRU
    seq_bottom = top_y + box_h
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, 
                                    Inches(6.0), seq_bottom + Inches(0.1), 
                                    Inches(0.5), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = DARK_BLUE
    arrow.line.fill.background()

    # GRU box
    gru_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                      Inches(3.5), seq_bottom + Inches(0.7), 
                                      Inches(5.5), Inches(0.7))
    gru_box.fill.solid()
    gru_box.fill.fore_color.rgb = ACCENT_RED
    gru_box.line.color.rgb = DARK_BLUE
    gru_box.line.width = Pt(2)
    tf = gru_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "GRU (hidden=32) processes 8 steps sequentially → FC(24) → FC(3)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Explanation below
    txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.5), Inches(2.2))
    tf = txBox3.text_frame
    tf.word_wrap = True
    explanations = [
        "• 8 consecutive feature vectors, each 200 ms apart → 1.6 seconds of temporal context",
        "• Each vector summarizes a 5-second window → total signal coverage: ~6.4 seconds",
        "• GRU hidden state (32 units) carries memory from step t-7 through step t",
        "• At step t, GRU has seen the full trajectory of how features evolved over 1.6s",
        "• Only the final hidden state (at step t) is passed to classification layers",
        "• This lets the model detect TRANSITIONS (blood→clot, clot→wall) not just static states",
    ]
    for i, txt in enumerate(explanations):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = txt
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)
    return slide

add_sequence_slide()

# Slide 12: GRU Hidden State — What the 32 Units Learn
def add_gru_hidden_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "GRU Hidden State — How 32 Units Encode Temporal Patterns"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left side: GRU cell diagram (simplified)
    # Title for left
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.5), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "GRU Cell (at each of 8 time steps):"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY

    # Simplified GRU equations box
    eq_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                     Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5))
    eq_box.fill.solid()
    eq_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    eq_box.line.color.rgb = MEDIUM_BLUE
    eq_box.line.width = Pt(2)
    tf = eq_box.text_frame
    tf.word_wrap = True
    equations = [
        "Update gate:  z = σ(W_z · [h_{t-1}, x_t])",
        "Reset gate:    r = σ(W_r · [h_{t-1}, x_t])",
        "Candidate:     ĥ = tanh(W · [r ⊙ h_{t-1}, x_t])",
        "New hidden:    h_t = (1-z) ⊙ h_{t-1} + z ⊙ ĥ",
        "",
        "x_t = 21 features (input at step t)",
        "h_t = 32 hidden units (carried to step t+1)",
    ]
    for i, eq in enumerate(equations):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = eq
        p.font.size = Pt(12)
        p.font.bold = (i < 4)
        p.font.color.rgb = DARK_BLUE if i < 4 else DARK_GRAY
        p.space_after = Pt(2)

    # Right side: What the hidden state captures
    tb2 = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.5))
    tf = tb2.text_frame
    p = tf.paragraphs[0]
    p.text = "What the 32 Hidden Units Learn:"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY

    concepts = [
        ("Update gate (z)", "How much old memory to keep vs. replace", MEDIUM_BLUE),
        ("Reset gate (r)", "Which past info is relevant to current input", MEDIUM_BLUE),
        ("Hidden state (h)", "Compressed temporal summary of the signal", ACCENT_RED),
    ]
    
    for i, (title, desc, color) in enumerate(concepts):
        y = Inches(2.1) + i * Inches(0.9)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                      Inches(6.8), y, Inches(5.8), Inches(0.75))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"{title}: {desc}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT

    # Bottom: Key insight
    txBox4 = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.3))
    tf = txBox4.text_frame
    tf.word_wrap = True
    insights = [
        "Key architectural choices:",
        "  • hidden_size=32: Just 32 numbers encode all temporal context — compact but sufficient",
        "  • GRU (not LSTM): Simpler (2 gates vs 3), fewer parameters, trains faster on small data",
        "  • Sequence=8 × 200ms = 1.6s: Captures tissue transitions without excessive delay",
        "  • Each 5s window already has temporal info → GRU adds inter-window evolution",
        "  • Total parameters: ~5,000 (GRU: ~3.5K + FC: ~1.5K) → resistant to overfitting",
    ]
    for i, txt in enumerate(insights):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = txt
        p.font.size = Pt(14)
        p.font.bold = (i == 0)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)
    return slide

add_gru_hidden_slide()

# Slide 13: Data Selection
add_content_slide("Data Selection Strategy", [
    "Filter: clot_mean_R < wall_mean_R (normal polarity studies only)",
    "Rationale: Absolute-level features (mean, max, range) require consistent polarity",
    "~45% of studies are 'inverted' (clot R > wall R) — excluded for now",
    "",
    "Current best: 34 training studies (April 6 data)",
    "Fixed test set: 8 studies held constant across all experiments",
    "",
    "Key finding: More data doesn't always improve results",
    "   • Quality and consistency matter more than quantity at this scale",
    "   • Best performance at 34 studies, worse at 41 studies",
])

# Save
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 
                        'LiveClotDetector_Architecture_v2.pptx')
prs.save(out_path)
print(f"Saved: {out_path}")
