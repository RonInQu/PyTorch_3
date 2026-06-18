"""Create visual-heavy PowerPoint for electrode impedance modeling proof of concept.

Uses .png figures from OpenFoam directory and SimResults.txt data.
Emphasizes visuals, capability demonstration, and proof-of-concept status.
Includes 2D model inaccuracies and COMSOL 3D value proposition.
"""

import os
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IMG_DIR = r"c:\Users\RonaldKurnik\OneDrive - Inquis Medical\Documents\2026\OpenFoam"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
NAVY = RGBColor(0x12, 0x2B, 0x45)
BLUE = RGBColor(0x1F, 0x77, 0xB4)
GREEN = RGBColor(0x2D, 0x8B, 0x4E)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED = RGBColor(0xC0, 0x39, 0x2B)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0xF2, 0xF5, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREEN = RGBColor(0x1A, 0x6B, 0x3A)
AMBER = RGBColor(0xF3, 0x9C, 0x12)


def img(name):
    return os.path.join(IMG_DIR, name)


def title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                  Inches(13.333), Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.12), Inches(12.5), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.35), Inches(0.58), Inches(12.5), Inches(0.3))
        sp = st.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(13)
        sp.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)


def bullets(slide, left, top, width, height, items, sz=14, color=GRAY, bold_first=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        if bold_first and i == 0:
            p.font.bold = True


def caption(slide, left, top, width, text, sz=10, color=GRAY):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER


def add_table(slide, left, top, width, height, headers, rows, hdr_color=BLUE, font_sz=10):
    ts = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    t = ts.table
    for c, h in enumerate(headers):
        cell = t.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = hdr_color
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(font_sz)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(rows, 1):
        for c, v in enumerate(row):
            cell = t.cell(r, c)
            cell.text = str(v)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(font_sz - 1)
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT


def poc_badge(slide, left=10.5, top=0.15):
    """Add a 'PROOF OF CONCEPT' badge."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                  Inches(2.5), Inches(0.35))
    box.fill.solid()
    box.fill.fore_color.rgb = AMBER
    box.line.fill.background()
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "PROOF OF CONCEPT"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Bioimpedance Electrode Modeling — Proof of Concept",
          "2D quasi-static FEA + Cole-Cole frequency analysis | Inquis Medical")
poc_badge(s)

# Two images side-by-side as visual hook
s.shapes.add_picture(img("Electric Potential.png"), Inches(0.5), Inches(1.3), Inches(6.2))
s.shapes.add_picture(img("E-Field Lines + Equipotentials.png"), Inches(6.8), Inches(1.3), Inches(6.2))

bullets(s, 0.8, 5.8, 12.0, 1.4, [
    "Objective: Explore whether multi-frequency impedance can distinguish clot from vessel wall.",
    "Status: 2D proof-of-concept model complete. Results are directionally correct but not yet quantitatively accurate.",
    "Geometry: 0.78 mm x 2.1 mm electrodes, 5.48 mm spacing, 50 kHz excitation.",
], sz=15)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — Geometry & Mesh
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "2D Model Setup: Geometry, Mesh, and Boundary Conditions")
poc_badge(s)

s.shapes.add_picture(img("2D Geometry.png"), Inches(0.3), Inches(1.15), Inches(6.4))
s.shapes.add_picture(img("Mesh.png"), Inches(6.8), Inches(1.15), Inches(6.2))

caption(s, 0.3, 4.6, 6.4, "Domain geometry with electrode edge labels")
caption(s, 6.8, 4.6, 6.2, "Quadratic FE mesh: 3750 nodes, 1815 elements")

bullets(s, 0.5, 5.0, 12.5, 2.2, [
    "Laplace equation solved: div(sigma * grad(V)) = 0 with Dirichlet BCs on electrodes (+/-1.5 V).",
    "All non-electrode boundaries: Neumann (insulating catheter surface or far-field).",
    "Cell constant K = 2812 m^-1 derived from numerical current integration.",
], sz=14)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — Field solution
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Electric Potential and Field Lines")
poc_badge(s)

s.shapes.add_picture(img("Electric Potential.png"), Inches(0.3), Inches(1.15), Inches(6.3))
s.shapes.add_picture(img("E-Field Lines + Equipotentials.png"), Inches(6.7), Inches(1.15), Inches(6.3))

caption(s, 0.3, 5.0, 6.3, "Potential contours (V range: -1.5 to +1.5 V)")
caption(s, 6.7, 5.0, 6.3, "E-field streamlines (red) perpendicular to equipotentials (blue)")

bullets(s, 0.5, 5.4, 12.5, 1.8, [
    "Field is concentrated between electrodes near the catheter surface — this is the sensing zone.",
    "Far-field decays rapidly; most sensitivity is within ~5 mm of the surface.",
], sz=14)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Impedance bar chart + frequency sweep
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Impedance by Material and Frequency Sweep", "Cole-Cole dispersion model calibrated to bench observations")
poc_badge(s)

s.shapes.add_picture(img("Impedance Comparison.png"), Inches(0.3), Inches(1.15), Inches(6.3))
s.shapes.add_picture(img("Impedance vs Frequency.png"), Inches(6.7), Inches(1.15), Inches(6.3))

caption(s, 0.3, 5.1, 6.3, "50 kHz snapshot: Blood 800, Clot 3486, Wall 1821 Ohm")
caption(s, 6.7, 5.1, 6.3, "Full sweep 1-100 kHz: magnitude (top) and phase (bottom)")

bullets(s, 0.5, 5.5, 12.5, 1.6, [
    "Clot/Blood = 4.36x,  Wall/Blood = 2.28x  at 50 kHz.   Saline is flat (no dispersion).",
    "All tissues converge at high frequency as dispersion saturates.",
], sz=14)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — Conductivity + discrimination
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Conductivity Dispersion and Clot vs Wall Discrimination")
poc_badge(s)

s.shapes.add_picture(img("Conductivity vs Frequency.png"), Inches(0.3), Inches(1.15), Inches(6.3))
s.shapes.add_picture(img("Clot vs Wall Frequency Discrimination.png"), Inches(6.7), Inches(1.15), Inches(6.3))

caption(s, 0.3, 5.1, 6.3, "Effective conductivity (top) and impedance ratio to blood (bottom)")
caption(s, 6.7, 5.1, 6.3, "Clot/Wall ratio with Monte Carlo uncertainty and conservative bands")

bullets(s, 0.5, 5.5, 12.5, 1.6, [
    "Clot/Wall ratio ranges from ~2.4x at low frequency to ~1.9x at 100 kHz.",
    "Monte Carlo p10-p90 band stays above 1 across the range — separation is robust under parameter uncertainty.",
], sz=14)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — Frequency discrimination visuals (replaces tables)
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Frequency Discrimination Summary", "Clot vs Wall impedance, ratio, contrast, and phase across frequency")
poc_badge(s)

s.shapes.add_picture(img("Frequency Discrimination Visual.png"), Inches(0.2), Inches(1.05), Inches(13.0))


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — 3-frequency feature comparison
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "3-Frequency Feature Set: 5, 50, 100 kHz", "Magnitude ratios and phase deltas normalized to 50 kHz anchor")
poc_badge(s)

s.shapes.add_picture(img("3-Frequency Feature Comparison.png"), Inches(0.3), Inches(1.1), Inches(9.5))

bullets(s, 10.0, 1.3, 3.1, 5.5, [
    "Recommended set:",
    "  5, 50, 100 kHz",
    "",
    "50 kHz = current HW",
    "5 kHz = max contrast",
    "100 kHz = spectral shape",
    "",
    "Best single extra freq:",
    "  ~5 kHz",
    "  Ratio = 2.11",
    "  dPhase = -7.2 deg",
    "  MC p10 = 1.50",
    "  MC p90 = 3.14",
], sz=12)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — E-Field Penetration Depth
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "E-Field Penetration Depth", "How far does the sensing field extend into tissue?")
poc_badge(s)

s.shapes.add_picture(img("E-Field Penetration Depth.png"), Inches(0.3), Inches(1.1), Inches(8.5))

bullets(s, 9.0, 1.2, 4.1, 5.5, [
    "Penetration from surface:",
    "",
    "Midpoint (z = 0):",
    "  1/e (37%):   3.05 mm",
    "  10%:           7.84 mm",
    "",
    "Electrode 1 center:",
    "  1/e (37%):   2.67 mm",
    "  10%:           7.53 mm",
    "",
    "Field decays rapidly;",
    "most sensitivity is within",
    "the first 3 mm of surface.",
], sz=13, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — Joule Heating / Power Density
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Joule Heating: Temperature Rise vs Distance from Electrode",
          "Adiabatic worst-case at 50 kHz, +/-1.5 V   |   IEC 60601 limit: dT < 2 deg C")
poc_badge(s)

s.shapes.add_picture(img("Temperature Rise vs Distance.png"), Inches(0.2), Inches(1.05), Inches(9.0))

bullets(s, 9.4, 1.2, 3.7, 5.5, [
    "Key observations:",
    "",
    "Saline (highest sigma):",
    "  Exceeds 2 C at surface",
    "  for both 1s and 10s.",
    "",
    "Blood:",
    "  Exceeds 2 C at surface",
    "  for 10s. Safe at >= 1 mm.",
    "",
    "Clot and Wall:",
    "  Well below 2 C at all",
    "  distances and durations.",
    "",
    "All materials safe at",
    ">= 2 mm depth.",
    "",
    "ADIABATIC (no cooling).",
    "Real convective cooling",
    "reduces actual rise.",
], sz=11, color=RED)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — Sensing Depth (cumulative energy)
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Sensing Depth: Where Does the Signal Come From?",
          "Cumulative fraction of total dissipated power vs distance from catheter surface")
poc_badge(s)

s.shapes.add_picture(img("Sensing Depth.png"), Inches(0.3), Inches(1.1), Inches(7.5))

bullets(s, 8.0, 1.3, 5.0, 5.0, [
    "Cumulative sensing energy:",
    "",
    "  50%  within  0.78 mm",
    "  80%  within  2.12 mm",
    "  95%  within  5.53 mm",
    "",
    "Implication:",
    "Half the signal comes from",
    "< 1 mm of the surface.",
    "",
    "The device is primarily a",
    "surface-contact sensor.",
    "",
    "Clot or wall must be in",
    "direct electrode proximity",
    "to dominate the measurement.",
], sz=14, color=NAVY)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — 2D Model Inaccuracies
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Inaccuracies of the Current 2D Model", "Why this is a proof of concept, not a design tool (yet)")
poc_badge(s)

# Left column
bullets(s, 0.5, 1.2, 6.0, 5.8, [
    "1. 2D vs 3D field geometry",
    "    Assumes infinite electrode width. Real 1 mm wide",
    "    electrodes have lateral field divergence. 2D",
    "    overestimates coupling, underestimates impedance.",
    "",
    "2. Flat surface approximation",
    "    Catheter is cylindrical (~8F = 2.7 mm OD).",
    "    Curved surface changes field radiation pattern.",
    "",
    "3. No device geometry",
    "    Real catheter has opening/window where electrodes",
    "    sit recessed. Cavity focuses the near-field.",
    "",
    "4. Rectangular domain vs cylindrical artery",
    "    PA is ~25 mm diameter. Field wraps differently",
    "    than in a rectangular half-space.",
], sz=14, color=RED)

# Right column
bullets(s, 6.7, 1.2, 6.3, 5.8, [
    "5. Homogeneous medium",
    "    One uniform material per solve. Real case has",
    "    partial clot contact creating asymmetric problem.",
    "",
    "6. No electrode-electrolyte interface",
    "    Real electrodes have double-layer capacitance",
    "    (CPE) that affects measured impedance at 50 kHz.",
    "",
    "7. No frequency dependence in field solve",
    "    Conductivity solve is purely real (DC-like).",
    "    Displacement current term is small (~2%) but",
    "    matters for low-conductivity wall tissue.",
    "",
    "Biggest single error: flat/infinite-width assumption",
    "likely underestimates impedance by 5-10x vs real 3D.",
], sz=14, color=RED)

bullets(s, 0.5, 6.6, 12.5, 0.6, [
    "Bottom line: 2D model shows correct qualitative field patterns and relative impedance ratios. "
    "Absolute values require the real 3D geometry.",
], sz=14, color=NAVY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 — What COMSOL 3D adds (capability/impact table)
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "What COMSOL 3D Simulation Adds", "From proof-of-concept to design tool")

add_table(s, 0.5, 1.2, 12.3, 5.0,
          ["Capability", "Impact"],
          [
              ["Import actual STEP geometry", "Correct cavity/electrode shape, realistic near-field"],
              ["3D solution", "Captures lateral spreading, true 3D current paths"],
              ["Multi-domain materials", "Model partial clot contact on one electrode only"],
              ["Complex-valued sigma*", "Proper phase information for impedance spectroscopy"],
              ["Cylindrical artery boundary", "Realistic far-field, correct current return paths"],
              ["Parametric sweeps", "Vary clot coverage %, thickness, position"],
              ["Contact impedance BC", "Model electrode-electrolyte double layer"],
              ["Frequency sweep", "Full impedance spectrum (not just one frequency)"],
              ["Adaptive meshing", "Fine mesh near electrode edges where E-field is singular"],
          ],
          hdr_color=DARK_GREEN, font_sz=14)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 — 10 Reasons for COMSOL
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "10 Reasons COMSOL Benefits Electrode Design Optimization")

# Left 5
bullets(s, 0.5, 1.15, 6.2, 5.7, [
    "1. Import actual CAD geometry",
    "   Real STEP files: cavity, recess, curved surfaces.",
    "   Design iterations in software before prototypes.",
    "",
    "2. Optimize electrode size, shape, spacing",
    "   Parametric sweeps identify geometry that maximizes",
    "   clot-wall contrast. Fewer physical prototyping cycles.",
    "",
    "3. Simulate partial coverage scenarios",
    "   Model 25/50/75% clot contact. Maps directly to",
    "   clinical sensitivity: how much clot is detectable?",
    "",
    "4. Multi-material domains",
    "   Blood + clot + wall simultaneously. This is the",
    "   real clinical scenario (heterogeneous problem).",
    "",
    "5. Frequency sweep optimization",
    "   Determine optimal multi-freq combo. Directly",
    "   informs firmware/ASIC design decisions.",
], sz=14, color=DARK_GREEN)

# Right 5
bullets(s, 6.7, 1.15, 6.3, 5.7, [
    "6. Predict absolute impedance values",
    "   Match simulation to bench and in-vivo data.",
    "   Predictive design before cutting metal.",
    "",
    "7. Sensing depth and field penetration",
    "   Does device sense only surface contact, or clot",
    "   behind thin blood layer? Critical for false negatives.",
    "",
    "8. Current density and safety analysis",
    "   Map heating at electrode edges. Required for",
    "   IEC 60601 compliance justification.",
    "",
    "9. Electrode material and coating effects",
    "   Model contact impedance for Pt, Au, PtIr, polymers.",
    "   Evaluate surface treatments vs sensitivity.",
    "",
    "10. Regulatory and IP documentation",
    "   Quantitative justification for 510(k)/De Novo.",
    "   Strengthens patent claims with reduction-to-practice.",
], sz=14, color=DARK_GREEN)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 — Without vs With COMSOL comparison
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Without Simulation vs With COMSOL", "Quantitative impact on development efficiency")

add_table(s, 1.5, 1.5, 10.3, 3.5,
          ["Without Simulation", "With COMSOL"],
          [
              ["5-10 physical prototypes per design iteration", "50+ virtual designs per day"],
              ["Bench testing reveals problems late", "Problems identified before fabrication"],
              ["\"It works but we don't know why\"", "Quantitative understanding of sensing mechanism"],
              ["Empirical frequency selection", "Optimized frequency backed by physics"],
              ["Regulatory asks \"how do you know it's safe?\"", "Simulation report demonstrates compliance"],
          ],
          hdr_color=NAVY, font_sz=14)

bullets(s, 1.5, 5.5, 10.3, 1.5, [
    "COMSOL transforms the design process from trial-and-error prototyping to physics-driven optimization.",
    "Simulation evidence directly supports regulatory submissions and IP filings.",
], sz=15, color=NAVY)





# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 — Next steps
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Recommended Next Steps")

bullets(s, 0.8, 1.4, 11.8, 5.5, [
    "1.  Lock initial 3-frequency firmware set: 5, 50, 100 kHz.",
    "",
    "2.  Capture phase-stable benchtop data for clot/wall mixtures and partial-contact cases.",
    "",
    "3.  Procure COMSOL AC/DC module license.",
    "",
    "4.  Build first 3D model with actual catheter STEP geometry + cylindrical vessel boundary.",
    "",
    "5.  Run parametric sweeps: electrode size/spacing, clot coverage %, wall thickness.",
    "",
    "6.  Run coupled electro-thermal analysis for worst-case operating conditions.",
    "",
    "7.  Validate simulation against benchtop measurements; iterate tissue parameters.",
    "",
    "8.  Use validated model to optimize electrode geometry before next hardware revision.",
], sz=16, color=GRAY)


# Footer on every slide
for sl in prs.slides:
    ft = sl.shapes.add_textbox(Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.25))
    p = ft.text_frame.paragraphs[0]
    p.text = f"Inquis Medical  |  Impedance Modeling Proof of Concept  |  {date.today().isoformat()}"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


out_path = r"c:\Users\RonaldKurnik\OneDrive - Inquis Medical\Documents\2026\PyTorch_3\Impedance_Frequency_Model_Summary_2026-06-17.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
