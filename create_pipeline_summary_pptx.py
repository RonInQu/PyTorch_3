"""Create a 2-slide PowerPoint summarizing the ML pipeline optimization journey."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MED_BLUE = RGBColor(0x2E, 0x5C, 0x8A)
LIGHT_BLUE = RGBColor(0x4A, 0x90, 0xD9)
GREEN = RGBColor(0x2D, 0x8B, 0x4E)
RED = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
GRAY = RGBColor(0x5D, 0x6D, 0x7E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 bold=False, color=DARK_BLUE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=11, color=DARK_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return tf


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Pipeline Architecture & Methods Tried
# ═══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

# Title bar
shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                 Inches(13.33), Inches(0.9))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()
add_text_box(slide1, Inches(0.3), Inches(0.15), Inches(12), Inches(0.7),
             "ML Clot Detection Pipeline — Methods & Architecture Exploration",
             font_size=22, bold=True, color=WHITE)

# Subtitle
add_text_box(slide1, Inches(0.3), Inches(0.95), Inches(12), Inches(0.4),
             "Inquis Medical — Resistance-Based Tissue Classification (Blood / Clot / Wall)",
             font_size=13, color=GRAY)

# ── Left column: Pipeline Architecture ──
add_text_box(slide1, Inches(0.3), Inches(1.5), Inches(4), Inches(0.4),
             "Final Pipeline Architecture", font_size=14, bold=True, color=MED_BLUE)

pipeline_steps = [
    "1. LabelingWithDuration — GT labels + 7s min filter",
    "2. Feature Extraction — 21 clot/wall-focused features",
    "     (5s window, 150 Hz, stride=30 samples)",
    "3. StandardScaler — per-feature normalization",
    "4. GRU(32) → FC(24, ReLU) → FC(3) classifier",
    "     (SEQ_LEN=8, batch=1024, seed=456)",
    "5. Inference: Temperature=1.5, EMA posterior",
    "     + DA override logic (confidence thresholds)",
]
add_bullet_list(slide1, Inches(0.3), Inches(1.9), Inches(4.5), Inches(3.5),
                pipeline_steps, font_size=11)

# ── Right column: Methods Tried ──
add_text_box(slide1, Inches(5.0), Inches(1.5), Inches(4), Inches(0.4),
             "Methods Explored", font_size=14, bold=True, color=MED_BLUE)

methods = [
    "DATA STRATEGIES:",
    "  • Cherry-picked subsets (24/31/34/41/85 studies)",
    "  • Polarity filtering (normal vs inverted cases)",
    "  • Duration filter: 0s, 3s, 7s min event length",
    "  • Label smoothing (CrossEntropy variant)",
    "",
    "FEATURE ENGINEERING (10+ sets tested):",
    "  • Baseline 40 statistical features",
    "  • +Hjorth, +CoeffVar, +Plateau, +Settling",
    "  • Short-timescale slopes (0.1–0.6s)",
    "  • Rise-shape morphology (7 amplitude-invariant)",
    "  • Texture RMS (bandpass 5–50 Hz)",
    "  • Pulse/cardiac features",
    "  • Division normalization (R/baseline)",
    "",
    "ARCHITECTURE VARIANTS:",
    "  • Single GRU(32) + FC head ← WINNER",
    "  • Stacked GRU1(32)→GRU2(24) (4 variants)",
    "  • CNN-1D (scattering transform)",
    "  • Denoised input experiments",
    "  • SEQ_LEN: 1, 8, 16",
]
add_bullet_list(slide1, Inches(5.0), Inches(1.9), Inches(4.3), Inches(5.2),
                methods, font_size=10)

# ── Far right: Key Decisions ──
add_text_box(slide1, Inches(9.5), Inches(1.5), Inches(3.5), Inches(0.4),
             "Key Decisions & Lessons", font_size=14, bold=True, color=MED_BLUE)

lessons = [
    "✓ 7s duration filter eliminates short",
    "   spikes → allows ALL 85 studies",
    "",
    "✓ Single GRU > Stacked GRU for",
    "   SEQ_LEN=8 (too few timesteps)",
    "",
    "✓ 21 clot/wall-focused features",
    "   outperform all 40+ variants",
    "",
    "✓ Absolute R level is key signal;",
    "   normalization destroys it",
    "",
    "✓ Mixing inverted polarity is",
    "   catastrophic (features contradict)",
    "",
    "✗ Feature ceiling reached —",
    "   orthogonal signals needed",
    "   (pressure, multi-frequency)",
]
add_bullet_list(slide1, Inches(9.5), Inches(1.9), Inches(3.6), Inches(5.0),
                lessons, font_size=10)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: Results Comparison Table
# ═══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

# Title bar
shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                 Inches(13.33), Inches(0.9))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()
add_text_box(slide2, Inches(0.3), Inches(0.15), Inches(12), Inches(0.7),
             "Results Summary — Net Benefit Progression (Same 8 Test Studies)",
             font_size=22, bold=True, color=WHITE)

# ── Main results table ──
from pptx.util import Cm

# Table data
headers = ["Experiment", "Train", "Config", "F1-macro", "Override\nPrecision",
           "Harmful", "Net Benefit"]
rows = [
    ["Batch 1 cherry-pick (Apr 28)", "24", "no filter, meanR<", "0.632", "0.532", "—", "+47,308"],
    ["Batch 2 mixed polarity", "47", "normal+inverted", "—", "—", "—", "-18,234"],
    ["Batch 2 normal-only", "31", "no filter, meanR<", "0.615", "—", "—", "+28,832"],
    ["Batch 3 no filter", "34", "no duration filter", "0.650", "0.574", "17,413", "+56,904"],
    ["Batch 3 + 3s filter", "34", "MIN_DUR=3s", "0.654", "0.573", "16,096", "+61,517"],
    ["Batch 3 + 7s filter", "34", "MIN_DUR=7s", "0.643", "0.617", "9,609", "+63,737"],
    ["Batch 4 + 7s filter", "41", "meanR<, 7s", "0.641", "0.565", "15,195", "+65,907"],
    ["ALL data + 7s ★ BEST", "85", "7s, no meanR filter", "0.650", "0.545", "11,938", "+75,420"],
    ["Stacked GRU (best)", "85", "GRU×2 + FC/ReLU/FC", "0.631", "0.449", "22,482", "+40,717"],
    ["Division normalization", "85", "R/baseline features", "—", "—", "—", "+6,973"],
    ["SEQ_LEN=1", "85", "no temporal context", "0.635", "0.478", "25,994", "+56,354"],
    ["Expanded test (14 studies)", "85", "harder test set", "0.649", "0.434", "24,990", "+17,595"],
]

n_rows = len(rows) + 1
n_cols = len(headers)
table_shape = slide2.shapes.add_table(n_rows, n_cols,
                                       Inches(0.3), Inches(1.2),
                                       Inches(12.7), Inches(4.5))
table = table_shape.table

# Set column widths
col_widths = [Inches(3.2), Inches(0.7), Inches(2.2), Inches(1.2),
              Inches(1.3), Inches(1.2), Inches(1.5)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

# Header row
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = MED_BLUE
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# Data rows
for r, row_data in enumerate(rows):
    for c, val in enumerate(row_data):
        cell = table.cell(r + 1, c)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(9)
        p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Highlight best row
        if "BEST" in row_data[0]:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xD5, 0xF5, 0xE3)
            p.font.bold = True
        # Highlight failures
        elif "mixed polarity" in row_data[0] or "normalization" in row_data[0]:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFA, 0xDB, 0xD8)

# ── Bottom summary ──
add_text_box(slide2, Inches(0.3), Inches(5.9), Inches(6.5), Inches(0.4),
             "Production Model: 2026-05-04_162814 — 85 studies, 7s filter, seed 456",
             font_size=12, bold=True, color=GREEN)

summary_text = [
    "Fixed test set: 33CFB812, 819421BC, 847A1E3F, 8ECEADA6, CENT0008, DD2DFAF4, F427536B, SUMM0127",
    "Architecture: GRU(input=21, hidden=32) → FC(32→24, ReLU) → FC(24→3, softmax)",
    "Net benefit = Correct overrides − Harmful overrides    |    Override precision target: >0.85",
]
add_bullet_list(slide2, Inches(0.3), Inches(6.3), Inches(12.5), Inches(1.1),
                summary_text, font_size=10, color=GRAY)

# ── Right side annotation ──
add_text_box(slide2, Inches(7.5), Inches(5.9), Inches(5.5), Inches(0.4),
             "Next: Pressure data (orthogonal signal) to break feature ceiling",
             font_size=12, bold=True, color=ORANGE)

# Save
out_path = r"c:\Users\RonaldKurnik\OneDrive - Inquis Medical\Documents\2026\PyTorch_3\ML_Pipeline_Optimization_Summary.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
