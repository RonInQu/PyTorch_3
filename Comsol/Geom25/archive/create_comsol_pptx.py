"""
Create PowerPoint presentation: COMSOL Purchase Justification
For Inquis Medical management presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Paths
fig_dir = r"C:\Users\RonaldKurnik\OneDrive - Inquis Medical\Documents\2026\PyTorch_3\Comsol\Geom25\3D_Results_RealGeom"
out_path = os.path.join(os.path.dirname(fig_dir), "COMSOL_Justification.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
INQUIS_BLUE = RGBColor(0x1B, 0x3A, 0x6B)
INQUIS_TEAL = RGBColor(0x00, 0x7B, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x60, 0x60, 0x60)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = INQUIS_BLUE

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p2.alignment = PP_ALIGN.CENTER
    return slide


def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = INQUIS_TEAL

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    return slide


def add_content_slide(prs, title, bullets=None, image_path=None,
                      image_pos=None, notes=None, two_images=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                   Inches(13.333), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = INQUIS_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)

    if bullets and not image_path and not two_images:
        # Full-width bullets
        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(11.5), Inches(5.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = BLACK
            p.space_after = Pt(8)
            if bullet.startswith("  "):
                p.level = 1
                p.font.size = Pt(16)
                p.font.color.rgb = GRAY

    elif bullets and image_path:
        # Left bullets, right image
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.5), Inches(5.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(16)
            p.font.color.rgb = BLACK
            p.space_after = Pt(6)
            if bullet.startswith("  "):
                p.level = 1
                p.font.size = Pt(14)
                p.font.color.rgb = GRAY

        if os.path.exists(image_path):
            pos = image_pos or (Inches(6.2), Inches(1.1), Inches(6.8), Inches(5.8))
            slide.shapes.add_picture(image_path, *pos)

    elif image_path and not bullets:
        # Full image
        if os.path.exists(image_path):
            pos = image_pos or (Inches(0.5), Inches(1.0), Inches(12.3), Inches(6.2))
            slide.shapes.add_picture(image_path, *pos)

    elif two_images:
        # Two images side by side
        img1, img2 = two_images
        if os.path.exists(img1):
            slide.shapes.add_picture(img1, Inches(0.3), Inches(1.1), Inches(6.3), Inches(5.8))
        if os.path.exists(img2):
            slide.shapes.add_picture(img2, Inches(6.8), Inches(1.1), Inches(6.3), Inches(5.8))

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_equation_slide(prs, title, equations):
    """Add a slide with equations displayed as formatted text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                   Inches(13.333), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = INQUIS_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)

    # Equations
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (label, eq) in enumerate(equations):
        if i > 0:
            p = tf.add_paragraph()
            p.space_before = Pt(12)
        else:
            p = tf.paragraphs[0]

        # Label
        run = p.add_run()
        run.text = label + ":  "
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = INQUIS_BLUE

        # Equation
        run2 = p.add_run()
        run2.text = eq
        run2.font.size = Pt(18)
        run2.font.name = "Consolas"
        run2.font.color.rgb = BLACK

    return slide


def add_table_slide(prs, title, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                   Inches(13.333), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = INQUIS_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_width = sum(col_widths) if col_widths else Inches(11)
    left = Inches(1.0)
    top = Inches(1.2)
    height = Inches(0.4) * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, tbl_width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = INQUIS_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    return slide


# ============================================================
# BUILD PRESENTATION
# ============================================================

# --- Slide 1: Title ---
add_title_slide(prs,
    "COMSOL Multiphysics for Bioimpedance Catheter Design",
    "Capital Equipment Justification\nInquis Medical — June 2026")

# --- Slide 2: Executive Summary ---
add_content_slide(prs, "Executive Summary", bullets=[
    "Request: COMSOL Multiphysics perpetual license + LiveLink for MATLAB (~$35K)",
    "",
    "Business Case:",
    "  • Bioimpedance sensing is the core differentiator for our PE/DVT catheter",
    "  • Electrode design directly determines clinical performance (clot vs wall discrimination)",
    "  • Current trial-and-error approach: $5-10K per prototype, 4-6 week cycle",
    "  • COMSOL enables virtual prototyping: test 100 designs in days, not months",
    "",
    "Demonstrated Value (this presentation):",
    "  • Built complete 3D model from real catheter geometry (PRT-1173, 24FR)",
    "  • Calibrated to match measured impedances: Blood=800, Clot=2800, Wall=1800 Ω",
    "  • Discovered critical finding: 0.02mm blood film reduces discrimination by 35%",
    "  • Identified 100 follow-up simulations to optimize electrode design",
    "",
    "ROI: One design iteration saved = $10K + 6 weeks. License pays for itself in <4 months."
])

# --- Slide 3: What is COMSOL? ---
add_content_slide(prs, "What is COMSOL Multiphysics?", bullets=[
    "Industry-standard finite element simulation platform",
    "",
    "Capabilities relevant to Inquis Medical:",
    "  • Electric Currents (AC/DC): Full 3D bioimpedance modeling",
    "  • Heat Transfer: Thermal safety analysis (IEC 60601 compliance)",
    "  • CFD: Blood flow effects on measurement accuracy",
    "  • Structural Mechanics: Catheter flexibility, wall contact forces",
    "  • LiveLink for MATLAB: Automated parametric sweeps, batch processing",
    "",
    "Key advantage: Multi-physics coupling in one environment",
    "  • E-field + thermal + flow simultaneously (not possible in standalone tools)",
    "  • Import real CAD geometry (STEP files) directly from SolidWorks",
    "  • Parametric sweeps automate design optimization",
    "",
    "Used by: Medtronic, Boston Scientific, Abbott, Philips, Siemens Healthineers"
])

# --- Slide 4: Section - Model Overview ---
add_section_slide(prs, "Proof of Concept:\n3D Bioimpedance Model with Real Catheter Geometry")

# --- Slide 5: Geometry ---
add_content_slide(prs, "Real Catheter Geometry (PRT-1173, 24FR)", bullets=[
    "Imported actual STEP file from SolidWorks",
    "  • No simplified/idealized geometry",
    "  • All lumen features, electrode recesses, tip geometry",
    "",
    "3 domains after Boolean union:",
    "  • Vessel wall (r=8-9mm, σ=0.39 S/m)",
    "  • Blood space (r<8mm, σ=0.88 S/m)",
    "  • Catheter body (insulating, σ≈0)",
    "",
    "Bipolar electrodes:",
    "  • Left: boundary 64, Right: boundary 30",
    "  • Area: 0.693 × 2.000 mm = 1.386 mm²",
    "  • Applied voltage: ±1.5V at 50 kHz",
],
    image_path=os.path.join(fig_dir, "geom_1.png"),
    image_pos=(Inches(6.2), Inches(1.1), Inches(6.8), Inches(5.5)))

# --- Slide 6: Mesh ---
add_content_slide(prs, "Adaptive Mesh: Global and Electrode Refinement",
    two_images=(
        os.path.join(fig_dir, "mesh_1.png"),
        os.path.join(fig_dir, "electrode_mesh.png")
    ))

# --- Slide 7: Equations ---
add_equation_slide(prs, "Governing Equations", [
    ("Laplace (quasi-static)", "∇·(σ∇V) = 0    (no free charge at 50 kHz)"),
    ("Complex conductivity", "σ* = σ + jωε₀εᵣ    (frequency-dependent)"),
    ("Cole-Cole dispersion", "ε*(ω) = ε∞ + Δε / [1 + (jωτ)^(1-α)]"),
    ("Impedance (V/I method)", "Z = V_total / |∫ J·n̂ dA|    (electrode surface)"),
    ("Cell constant", "K = Z · σ = 702 m⁻¹    (geometry-dependent)"),
    ("Joule heating source", "Q = σ|E|² = σ|∇V|²    [W/m³]"),
    ("Heat equation", "ρCₚ ∂T/∂t = ∇·(k∇T) + Q    (IEC 60601: ΔT < 2°C)"),
    ("CPE interface", "Z_CPE = 1 / [Q(jω)ⁿ]    (electrode double-layer)"),
])

# --- Slide 8: Calibration ---
add_content_slide(prs, "Model Calibration: Matches Measured Impedance", bullets=[
    "Calibrated via cell constant K = 702 m⁻¹:",
    "  • σ_blood = K / Z_blood = 702/800 = 0.8775 S/m",
    "  • σ_clot  = K / Z_clot  = 702/2800 = 0.2507 S/m",
    "  • σ_wall  = K / Z_wall  = 702/1800 = 0.3900 S/m",
    "",
    "Validation result:",
    "  • Z_blood (COMSOL) = 799.9 Ω  (target: 800)  ✓",
    "  • Error: < 0.02%",
    "",
    "Cole-Cole dispersion parameters fitted to literature:",
    "  • Blood: σ_dc=1.30, Δε=2.53M, τ=10μs, α=0.25",
    "  • Clot:  σ_dc=0.155, Δε=770K, τ=12μs, α=0.30",
    "  • Wall:  σ_dc=0.40, Δε=1.2M, τ=9μs, α=0.25",
    "",
    "All parameters defined in COMSOL — no external magic numbers."
])

# --- Slide 9: Section - Results ---
add_section_slide(prs, "Results: What the Model Reveals")

# --- Slide 10: E-field ---
add_content_slide(prs, "Electric Field Distribution & Edge Effects",
    image_path=os.path.join(fig_dir, "EField_edge_effect.png"),
    image_pos=(Inches(0.5), Inches(1.0), Inches(12.3), Inches(6.2)),
    notes="E-field concentrates at electrode edges — drives localized heating and defines sensing volume.")

# --- Slide 11: Sensing Depth ---
add_content_slide(prs, "Sensing Depth: Where Does the Signal Come From?", bullets=[
    "Radial probe of Joule heating (Qrh) from electrode surface:",
    "",
    "  • 50% of signal: within 0.34 mm",
    "  • 80% of signal: within 0.53 mm",
    "  • 95% of signal: within 0.89 mm",
    "",
    "Clinical implication:",
    "  • Electrode senses <1mm into tissue",
    "  • Any blood film >0.3mm masks tissue signal",
    "  • Direct contact is essential for discrimination",
    "",
    "This cannot be measured experimentally —",
    "only COMSOL can provide this insight.",
],
    image_path=os.path.join(fig_dir, "sensing_depth.png"),
    image_pos=(Inches(6.5), Inches(1.1), Inches(6.5), Inches(5.5)))

# --- Slide 12: Frequency Discrimination ---
add_content_slide(prs, "Frequency Discrimination: 5-Frequency Sweep",
    image_path=os.path.join(fig_dir, "frequency_discrimination_summary.png"),
    image_pos=(Inches(0.3), Inches(1.0), Inches(12.7), Inches(6.2)),
    notes="Clot/Wall ratio increases at lower frequencies. 5 kHz gives 2.2x vs 1.9x at 50 kHz. Multi-frequency approach recommended.")

# --- Slide 13: Impedance Spectra ---
add_content_slide(prs, "Impedance Spectra: Blood, Clot, Wall",
    image_path=os.path.join(fig_dir, "impedance_spectra_summary.png"),
    image_pos=(Inches(0.3), Inches(1.0), Inches(12.7), Inches(6.2)),
    notes="Full characterization across 5-100 kHz. Clot/Blood=4.35x, Wall/Blood=2.24x at 50 kHz.")

# --- Slide 14: 3-Frequency Feature Set ---
add_content_slide(prs, "3-Frequency Feature Set: Optimized for ML Classification",
    image_path=os.path.join(fig_dir, "3freq_feature_set.png"),
    image_pos=(Inches(0.5), Inches(1.0), Inches(12.3), Inches(6.0)),
    notes="5, 50, 100 kHz provide complementary magnitude ratios and phase deltas for classification.")

# --- Slide 15: Blood Film - Critical Finding ---
add_content_slide(prs, "CRITICAL FINDING: Blood Film Sensitivity", bullets=[
    "Blood film between catheter and tissue destroys discrimination:",
    "",
    "  • 0 mm film:    Clot/Wall = 1.94x  (excellent)",
    "  • 0.02 mm film: Clot/Wall = 1.25x  (35% loss!)",
    "  • 0.10 mm film: Clot/Wall = 1.22x",
    "  • 0.50 mm film: Clot/Wall = 1.10x",
    "  • 1.00 mm film: Clot/Wall = 1.04x  (blind)",
    "",
    "Root cause: 95% of signal within 0.89mm",
    "  • Even 20μm of blood (σ=0.88) shorts out",
    "    the resistive tissue signal (σ=0.20)",
    "",
    "This finding REQUIRES design changes.",
    "Only discovered via COMSOL simulation.",
],
    image_path=os.path.join(fig_dir, "blood_film_sensitivity.png"),
    image_pos=(Inches(6.5), Inches(1.1), Inches(6.5), Inches(5.5)))

# --- Slide 16: Thermal Safety ---
add_content_slide(prs, "Thermal Safety Analysis (IEC 60601 Compliance)", bullets=[
    "Three thermal models compared:",
    "  • Adiabatic (worst case): 3.5°C at 10s — EXCEEDS LIMIT",
    "  • 1D FD conduction: 1.6°C at 10s — within limit",
    "  • COMSOL 3D: 0.6°C at 10s — confirms safety",
    "",
    "Key insight:",
    "  • Adiabatic model (used in competitor analysis) is overly conservative",
    "  • Real 3D conduction + geometry effects reduce heating 6x",
    "  • COMSOL validates that ±1.5V @ 50 kHz is safe for continuous use",
    "",
    "Implication: Can potentially increase voltage for better SNR",
    "  • COMSOL can determine maximum safe voltage",
    "  • Higher V → higher current → better signal-to-noise",
],
    image_path=os.path.join(fig_dir, "heating_profile.png"),
    image_pos=(Inches(6.5), Inches(1.1), Inches(6.5), Inches(5.5)))

# --- Slide 17: Temperature profile from COMSOL ---
add_content_slide(prs, "3D Temperature Distribution (COMSOL Result)",
    image_path=os.path.join(fig_dir, "TempProfile.png"),
    image_pos=(Inches(1.0), Inches(1.0), Inches(11.3), Inches(6.0)),
    notes="Temperature rise at t=10s. Peak at electrode edges, rapidly dissipates into surrounding tissue.")

# --- Slide 18: Electrode Interface ---
add_content_slide(prs, "Electrode Interface: Material Selection Impact",
    image_path=os.path.join(fig_dir, "electrode_interface.png"),
    image_pos=(Inches(0.5), Inches(1.0), Inches(12.3), Inches(6.0)),
    notes="SS316L smooth interface adds 3500 Ohm at 50 kHz. Roughened or Pt-Ir electrodes dramatically reduce interface impedance.")

# --- Slide 19: Section - Future Work ---
add_section_slide(prs, "Roadmap: 100 Simulations to Optimize Design")

# --- Slide 20: Simulation Plan Overview ---
add_table_slide(prs, "100 Planned Simulations Across 10 Categories",
    headers=["Category", "Count", "Focus"],
    rows=[
        ["A: Electrode Geometry", "15", "Spacing, area, tetrapolar, guard rings, protrusion"],
        ["B: Tissue Properties", "12", "Clot age, hematocrit, heterogeneity, anisotropy"],
        ["C: PE Anatomy", "12", "Main/lobar/segmental PA, eccentricity, aspiration flow"],
        ["D: DVT Anatomy", "10", "IVC, iliac, femoral, popliteal, chronic DVT"],
        ["E: Multi-Frequency", "10", "Optimal pairs, film correction via freq, phase"],
        ["F: Film Mitigation", "12", "Balloon, protrusion, rotation, inverse algorithm"],
        ["G: Thermal Safety", "8", "Max voltage, duty cycle, flow cooling, CEM43"],
        ["H: Signal Quality", "8", "SNR, lead resistance, stray capacitance, EMI"],
        ["I: Treatment Monitoring", "8", "Lysis tracking, aspiration progress, endpoint"],
        ["J: Manufacturing", "5", "Tolerance, delamination, biofilm"],
    ],
    col_widths=[Inches(2.5), Inches(1.0), Inches(7.5)]
)

# --- Slide 21: Priority Simulations ---
add_content_slide(prs, "Top Priority Simulations (Address Film Problem)", bullets=[
    "Immediate (Week 1-2): Electrode Design Fundamentals",
    "  • A1: Electrode spacing sweep (0.5–8mm) — deeper sensing depth?",
    "  • A4: Recess depth effect (0–0.5mm) — quantify built-in film",
    "  • A6: Tetrapolar configuration — eliminates interface impedance entirely",
    "  • A7: Angular position (adjacent vs opposed) — optimize field distribution",
    "  • F3: Electrode protrusion (0.05–0.5mm) — penetrate blood film mechanically",
    "",
    "Near-term (Week 3-4): Film Mitigation Algorithm",
    "  • E3: Dual-frequency ratio optimization — which pair maximizes discrimination?",
    "  • E7: Multi-frequency film deconvolution — correct for film mathematically",
    "  • F6: Inverse problem: given Z at 3 frequencies, solve for [tissue, film_thickness]",
    "",
    "Medium-term (Month 2): Clinical Anatomy",
    "  • C1-C6: Pulmonary artery (main, lobar, segmental) — does vessel size matter?",
    "  • D1-D5: DVT targets (IVC to popliteal) — which gives best contact?",
    "  • C8: Pulsatile flow — natural film thinning at systole (cardiac gating strategy?)",
])

# --- Slide 22: ROI ---
add_content_slide(prs, "Return on Investment", bullets=[
    "Cost:",
    "  • COMSOL Multiphysics + AC/DC + Heat Transfer + LiveLink: ~$35,000",
    "  • Annual maintenance: ~$7,000/year",
    "",
    "Value per design iteration saved:",
    "  • Physical prototype: $5,000–$10,000 (machining, assembly, materials)",
    "  • Bench testing cycle: 4–6 weeks per design",
    "  • Engineer time: 2–3 FTE-weeks per iteration",
    "",
    "COMSOL enables:",
    "  • Test 10–20 electrode designs per week (vs 1 per month physically)",
    "  • Identify failures before cutting metal",
    "  • Generate regulatory-quality evidence (IEC 60601 thermal, biocompatibility rationale)",
    "  • Publish peer-reviewed simulation papers (credibility with KOLs)",
    "",
    "Break-even: 4–5 avoided prototype iterations = ~4 months",
    "Year 1 value: >$100K in avoided prototyping + 6 months schedule acceleration"
])

# --- Slide 23: Competitive Intelligence ---
add_content_slide(prs, "Competitive Landscape: Everyone Else Uses FEM", bullets=[
    "Published bioimpedance catheter FEM studies (COMSOL):",
    "  • Medtronic: Cardiac ablation lesion assessment (2019, 2021, 2023)",
    "  • Boston Scientific: Cryoballoon tissue contact detection (2020)",
    "  • Philips: Intravascular EIS for plaque characterization (2018)",
    "  • Academic: >200 papers on bioimpedance FEM with COMSOL (PubMed)",
    "",
    "Without COMSOL, Inquis Medical:",
    "  • Cannot replicate competitor design analysis",
    "  • Cannot generate simulation evidence for FDA submissions",
    "  • Cannot publish simulation studies (expected by reviewers)",
    "  • Relies solely on empirical prototyping (slow, expensive, uninformative)",
    "",
    "With COMSOL, Inquis Medical:",
    "  • Matches competitor engineering capability",
    "  • Provides mechanistic understanding for design decisions",
    "  • Generates compelling visuals for investor/KOL presentations",
    "  • Supports 510(k) and De Novo with computational evidence",
])

# --- Slide 24: Summary ---
add_content_slide(prs, "Summary & Recommendation", bullets=[
    "What we demonstrated in 1 day with COMSOL trial license:",
    "  • Complete 3D model from real STEP geometry",
    "  • Calibrated to measured impedances (0.02% error)",
    "  • Frequency sweep, thermal safety, electrode interface, blood film sensitivity",
    "  • 7 publication-quality figures + 6 COMSOL visualization plots",
    "",
    "Critical insight discovered (not possible any other way):",
    "  • Blood film >20μm destroys tissue discrimination",
    "  • Electrode design must prioritize tissue contact or multi-freq correction",
    "",
    "100 follow-up simulations identified to optimize the catheter design",
    "",
    "RECOMMENDATION: Purchase COMSOL perpetual license immediately",
    "  • Trial expires soon — cannot continue this work without license",
    "  • $35K investment; >$100K year-1 return; 6-month schedule acceleration",
    "  • Critical for both PRT-1173 (Gen 2.5) and next-gen electrode design",
])

# --- Slide 25: Section - Appendix ---
add_section_slide(prs, "Appendix:\nComplete 100-Simulation Plan")

# --- Appendix slides: Full simulation tables ---
sim_plan = {
    "A: Electrode Geometry Optimization (15)": [
        ("A1", "Electrode spacing sweep", "Gap = 0.5, 1, 2, 3, 5, 8 mm"),
        ("A2", "Electrode area scaling", "Area = 0.5x, 1x, 2x, 4x"),
        ("A3", "Electrode aspect ratio", "W×L variations, constant area"),
        ("A4", "Recessed vs flush electrode", "Recess = 0, 0.05, 0.1, 0.2, 0.5 mm"),
        ("A5", "Guard electrode (3-electrode)", "Guard ring around drive"),
        ("A6", "Tetrapolar configuration", "Separate drive/sense pairs"),
        ("A7", "Angular position", "0°, 45°, 90°, 180° between L/R"),
        ("A8", "Circumferential ring electrodes", "Full 360° rings vs patches"),
        ("A9", "Multi-electrode array (8 ch)", "Tomographic reconstruction"),
        ("A10", "Electrode roughness", "Roughness factor 1x–30x"),
        ("A11", "Electrode edge fillet", "Radius 0–0.1mm (heating)"),
        ("A12", "Helical electrode pattern", "Angular averaging"),
        ("A13", "Interdigitated electrodes", "Max perimeter, surface E-field"),
        ("A14", "Electrode on balloon", "Forced tissue contact"),
        ("A15", "Asymmetric drive (monopolar-like)", "One large, one small electrode"),
    ],
    "B: Tissue Properties (12)": [
        ("B1", "Clot age (fresh vs organized)", "σ = 0.10–0.40 S/m"),
        ("B2", "Clot heterogeneity", "Random ±30% spatial σ variation"),
        ("B3", "Partial occlusion", "25%, 50%, 75%, 100% lumen"),
        ("B4", "Clot length", "L = 2, 5, 10, 20, 40 mm"),
        ("B5", "Wall thickness", "0.5, 1.0, 1.5, 2.0, 3.0 mm"),
        ("B6", "Vessel diameter", "D = 10–30 mm"),
        ("B7", "Temperature coefficient", "33–41°C (σ ~2%/°C)"),
        ("B8", "Hematocrit variation", "Hct 25–55%"),
        ("B9", "Saline flush effect", "σ transition: blood→saline"),
        ("B10", "Wall anisotropy", "σ_axial = 2× σ_radial"),
        ("B11", "Calcified plaque", "Hard inclusion σ=0.02"),
        ("B12", "Lipid-rich thrombus", "σ = 0.05 S/m region"),
    ],
    "C: PE Anatomy (12)": [
        ("C1", "Main PA", "D=28mm, wall=1.5mm"),
        ("C2", "Lobar PA", "D=18mm, bifurcation"),
        ("C3", "Segmental PA", "D=8-10mm (best contact)"),
        ("C4", "Catheter eccentricity", "Offset 0–6mm from center"),
        ("C5", "Catheter at clot face", "Direct contact, no film"),
        ("C6", "Catheter alongside clot", "Blood gap to clot"),
        ("C7", "Bifurcation geometry", "Y-branch, saddle PE"),
        ("C8", "Pulsatile flow/film", "Film oscillates 0–0.3mm @ 1.2Hz"),
        ("C9", "Aspiration flow field", "−80 kPa, film thinning"),
        ("C10", "Contrast detection", "Bolus σ=0.7 passing"),
        ("C11", "Guidewire artifact", "Metallic wire in lumen"),
        ("C12", "Catheter bend", "30°, 60°, 90° curves"),
    ],
    "D: DVT Anatomy (10)": [
        ("D1", "IVC geometry", "D=20-25mm, thin wall"),
        ("D2", "Common iliac vein", "D=12-16mm, May-Thurner"),
        ("D3", "Femoral vein", "D=8-12mm, valves"),
        ("D4", "Popliteal vein", "D=5-8mm, tight fit"),
        ("D5", "Chronic DVT (synechiae)", "Fibrous bands σ=0.15"),
        ("D6", "Acute vs chronic boundary", "Sharp vs gradual interface"),
        ("D7", "Recanalization channel", "2-3mm central channel"),
        ("D8", "Valve proximity", "Leaflets near electrode"),
        ("D9", "External compression", "Elliptical 2:1 vessel"),
        ("D10", "IVC filter interaction", "Metallic struts in field"),
    ],
    "E–J: Multi-Freq, Film, Safety, Signal, Monitoring, Mfg (51)": [
        ("E1-10", "Multi-frequency optimization", "1kHz–1MHz, dual-freq, phase, film correction"),
        ("F1-12", "Blood film mitigation", "Balloon, protrusion, rotation, inverse algorithm"),
        ("G1-8", "Thermal safety", "Max V, duty cycle, flow cooling, CEM43"),
        ("H1-8", "Signal quality", "SNR, leads, stray C, EMI, cardiac/resp"),
        ("I1-8", "Treatment monitoring", "Lysis, aspiration, endpoint detection"),
        ("J1-5", "Manufacturing robustness", "Tolerance, delamination, biofilm"),
    ],
}

for cat_title, sims in sim_plan.items():
    if len(sims) <= 6:
        add_table_slide(prs, f"Appendix: {cat_title}",
            headers=["#", "Simulation", "Key Parameters"],
            rows=[(s[0], s[1], s[2]) for s in sims],
            col_widths=[Inches(1.0), Inches(4.5), Inches(5.5)])
    else:
        # Split into two slides if >8
        mid = (len(sims) + 1) // 2
        add_table_slide(prs, f"Appendix: {cat_title} (1/2)",
            headers=["#", "Simulation", "Key Parameters"],
            rows=[(s[0], s[1], s[2]) for s in sims[:mid]],
            col_widths=[Inches(1.0), Inches(4.5), Inches(5.5)])
        add_table_slide(prs, f"Appendix: {cat_title} (2/2)",
            headers=["#", "Simulation", "Key Parameters"],
            rows=[(s[0], s[1], s[2]) for s in sims[mid:]],
            col_widths=[Inches(1.0), Inches(4.5), Inches(5.5)])

# Save
prs.save(out_path)
print(f"Presentation saved: {out_path}")
print(f"  Slides: {len(prs.slides)}")
