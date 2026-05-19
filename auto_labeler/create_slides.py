"""
Generate PowerPoint slides explaining the 1D U-Net auto-labeler architecture.
Run: python auto_labeler/create_slides.py
Output: auto_labeler/UNet_Architecture.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point
        tf.paragraphs[i].font.size = Pt(18)
    return slide

def add_architecture_diagram(prs):
    """Create a visual architecture diagram slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    slide.shapes.title.text = "1D U-Net Architecture"

    # Colors
    encoder_color = RGBColor(0x41, 0x69, 0xE1)  # Royal blue
    decoder_color = RGBColor(0x22, 0x8B, 0x22)  # Forest green
    bottleneck_color = RGBColor(0xFF, 0x8C, 0x00)  # Dark orange
    skip_color = RGBColor(0x80, 0x80, 0x80)  # Gray
    input_color = RGBColor(0x55, 0x55, 0x55)  # Dark gray

    # Encoder blocks (left side, descending)
    encoder_labels = [
        "Enc1: 32ch × 4096",
        "Enc2: 64ch × 2048",
        "Enc3: 128ch × 1024",
        "Enc4: 256ch × 512",
        "Enc5: 512ch × 256",
    ]

    # Decoder blocks (right side, ascending)
    decoder_labels = [
        "Dec5: 512ch × 256",
        "Dec4: 256ch × 512",
        "Dec3: 128ch × 1024",
        "Dec2: 64ch × 2048",
        "Dec1: 32ch × 4096",
    ]

    left_x = Inches(0.5)
    right_x = Inches(7.0)
    top_y = Inches(1.8)
    box_w = Inches(2.8)
    box_h = Inches(0.65)
    spacing = Inches(0.80)

    # Draw Input box
    input_y = top_y - Inches(0.9)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, input_y, box_w, Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = input_color
    tf = shape.text_frame
    tf.text = "Input: 1ch × 4096 (resistance)"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Draw encoder blocks
    for i, label in enumerate(encoder_labels):
        y = top_y + i * spacing
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = encoder_color
        tf = shape.text_frame
        tf.text = label
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Draw bottleneck
    bn_y = top_y + 5 * spacing
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), bn_y, Inches(3.0), box_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bottleneck_color
    tf = shape.text_frame
    tf.text = "Bottleneck: 1024ch × 128"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Draw decoder blocks
    for i, label in enumerate(decoder_labels):
        y = top_y + (4 - i) * spacing
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = decoder_color
        tf = shape.text_frame
        tf.text = label
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Draw skip connections (dashed arrows from encoder to decoder at same level)
    from pptx.oxml.ns import qn
    for i in range(5):
        enc_y = top_y + i * spacing + box_h / 2
        dec_y = top_y + (4 - i) * spacing + box_h / 2
        arrow_y = (enc_y + dec_y) / 2 if enc_y != dec_y else enc_y

        # Horizontal arrow from encoder right edge to decoder left edge
        arrow_left = left_x + box_w + Inches(0.1)
        arrow_right = right_x - Inches(0.1)
        arr_y = top_y + i * spacing + box_h / 2

        connector = slide.shapes.add_connector(
            1,  # straight connector
            arrow_left, arr_y,
            arrow_right, arr_y,
        )
        connector.line.color.rgb = skip_color
        connector.line.width = Pt(1.5)
        connector.line.dash_style = 2  # dash

        # Label the skip connection
        skip_label_x = Inches(3.8)
        skip_label_y = arr_y - Inches(0.15)
        txBox = slide.shapes.add_textbox(skip_label_x, skip_label_y, Inches(1.5), Inches(0.3))
        tf = txBox.text_frame
        tf.text = f"skip {i+1}"
        tf.paragraphs[0].font.size = Pt(8)
        tf.paragraphs[0].font.color.rgb = skip_color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Output
    out_y = top_y - Inches(0.9)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, out_y, box_w, Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x8B, 0x00, 0x00)
    tf = shape.text_frame
    tf.text = "Output: 3 classes × 4096"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ──
    add_title_slide(prs,
        "1D U-Net Auto-Labeler",
        "Offline Tissue Segmentation for Impedance Data\n"
        "Architecture Deep Dive\n"
        "May 2026"
    )

    # ── Slide 2: What is U-Net? ──
    add_content_slide(prs, "What is U-Net?", [
        "Originally designed for biomedical image segmentation (2015)",
        "Encoder-Decoder architecture with skip connections",
        "Encoder: extracts features at multiple scales (context)",
        "Decoder: reconstructs per-sample predictions (localization)",
        "Skip connections: preserve fine-grained details lost during downsampling",
        "Our version: adapted from 2D images → 1D time series",
        "Input: raw resistance signal (single channel, z-normalized)",
        "Output: per-sample class probability (blood/clot/wall)",
    ])

    # ── Slide 3: Architecture Overview ──
    add_content_slide(prs, "Architecture Overview", [
        "5-level encoder + bottleneck + 5-level decoder",
        "Input: 1 channel × 4096 samples (~25 seconds @ 167 Hz)",
        "Output: 3 classes × 4096 samples (one label per sample)",
        "Total parameters: 27.3 million",
        "Kernel size: 7 (captures ~42ms per convolution)",
        "Optional dropout (0.2) after each encoder + bottleneck",
        "Receptive field: ~500+ samples (~3 seconds) at bottleneck",
        "Filter progression: 32 → 64 → 128 → 256 → 512 → 1024 (bottleneck)",
    ])

    # ── Slide 4: Visual Diagram ──
    add_architecture_diagram(prs)

    # ── Slide 5: Encoder Path (Contracting) ──
    add_content_slide(prs, "Encoder Path — Feature Extraction", [
        "Each encoder level does:",
        "  1. Two Conv1d(kernel=7) + BatchNorm + ReLU",
        "  2. Dropout1d (optional, default 0.2 for regularization)",
        "  3. MaxPool1d(2) — halves the sequence length",
        "",
        "Level 1: 1→32 filters, length 4096 → 2048",
        "Level 2: 32→64 filters, length 2048 → 1024",
        "Level 3: 64→128 filters, length 1024 → 512",
        "Level 4: 128→256 filters, length 512 → 256",
        "Level 5: 256→512 filters, length 256 → 128",
        "",
        "Deeper = more abstract features, larger context window",
    ])

    # ── Slide 6: Bottleneck ──
    add_content_slide(prs, "Bottleneck — Maximum Context", [
        "512 → 1024 channels at length 128",
        "Each sample here 'sees' ~32× the original resolution",
        "With kernel=7, effective receptive field spans seconds of signal",
        "Dropout applied here too for regularization",
        "This is where the model captures long-range patterns:",
        "  • Sustained resistance elevations (clot events)",
        "  • Gradual transitions (wall contact)",
        "  • Overall signal dynamics",
        "",
        "The most compressed representation of the input",
    ])

    # ── Slide 7: Decoder Path (Expanding) ──
    add_content_slide(prs, "Decoder Path — Precise Localization", [
        "Each decoder level does:",
        "  1. ConvTranspose1d(stride=2) — doubles length",
        "  2. Concatenate with skip connection from encoder",
        "  3. Two Conv1d(kernel=7) + BatchNorm + ReLU",
        "",
        "Level 5: 1024→512, length 128 → 256  (+ skip from Enc5)",
        "Level 4: 512→256, length 256 → 512  (+ skip from Enc4)",
        "Level 3: 256→128, length 512 → 1024  (+ skip from Enc3)",
        "Level 2: 128→64, length 1024 → 2048  (+ skip from Enc2)",
        "Level 1: 64→32, length 2048 → 4096  (+ skip from Enc1)",
    ])

    # ── Slide 8: Skip Connections ──
    add_content_slide(prs, "Skip Connections — Why They Matter", [
        "Problem: downsampling loses spatial precision",
        "Solution: concatenate encoder features directly to decoder",
        "",
        "Encoder features provide: fine edges, local patterns, exact positions",
        "Decoder features provide: semantic context, what class it likely is",
        "Combined: precise boundaries + correct classification",
        "",
        "Without skips: output labels would be 'blurry' — poor boundary accuracy",
        "With skips: crisp transitions between blood/clot/wall segments",
    ])

    # ── Slide 9: Output Head ──
    add_content_slide(prs, "Output Head — Per-Sample Classification", [
        "1×1 Convolution: 32 channels → 3 channels",
        "Each output channel = logit for one class:",
        "  Channel 0: blood probability",
        "  Channel 1: clot probability",
        "  Channel 2: wall probability",
        "",
        "Training: CrossEntropyLoss with class weights [1.0, 13.0, 22.0]",
        "  (Also supports FocalLoss for harder mining — gamma=2.0)",
        "Inference: softmax → argmax per sample → predicted label",
        "",
        "Post-processing: minimum duration filter removes spurious short segments",
    ])

    # ── Slide 10: Key Numbers ──
    add_content_slide(prs, "Key Numbers", [
        "Parameters: 27,263,875 (27.3M)",
        "Chunk size: 4096 samples = 25 seconds @ 167 Hz",
        "Stride: 2048 (50% overlap between training chunks)",
        "Kernel size: 7 (each conv sees 7 consecutive samples = 42ms)",
        "Depth: 5 levels → 2^5 = 32× downsampling at bottleneck",
        "Dropout: 0.2 (optional, applied after encoder blocks + bottleneck)",
        "Classes: 3 (blood=0, clot=1, wall=2)",
        "Class weights: blood=1.0, clot=13.0, wall=22.0 (inverse frequency)",
        "Training studies: 85 | Validation: 12 (study-level split)",
        "Training time: ~10 min on A100 GPU (80 epochs, batch_size=64)",
    ])

    # ── Slide 11: Inference Pipeline ──
    add_content_slide(prs, "Inference Pipeline", [
        "1. Load full-length parquet (any duration)",
        "2. Z-score normalize per file: (R - mean) / std",
        "3. Slice into overlapping chunks (stride = chunk_size / 2)",
        "4. Run each chunk through U-Net → softmax probabilities",
        "5. Average probabilities in overlap regions",
        "6. Argmax → per-sample predicted label",
        "7. Post-process: remove segments < 6 seconds (spurious)",
        "8. Add 'predicted_label' column to parquet and save",
        "",
        "Result: labeled parquet matching original time resolution",
    ])

    # ── Slide 12: Results Summary ──
    add_content_slide(prs, "Results — 8 Test Studies (Production Model)", [
        "Best model: Run 1 (single-channel, CE loss, lr=1e-3, no dropout)",
        "",
        "Mean test F1 (macro): 0.688",
        "  Blood F1: 0.948 (excellent)",
        "  Wall F1:  0.803 (good)",
        "  Clot F1:  0.314 (moderate — known hard problem)",
        "",
        "Best individual: F427536B — F1=0.857",
        "Worst individual: DD2DFAF4 — F1=0.487",
        "",
        "Use case: draft labeler with human-in-the-loop review",
        "Multi-channel (5ch) + focal loss tested — didn't improve generalization",
    ])

    # Save
    out_path = os.path.join(os.path.dirname(__file__), "UNet_Architecture.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
